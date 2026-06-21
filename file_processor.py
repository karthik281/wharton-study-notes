"""Extract readable text from PDF, PPTX, and DOCX files."""

import logging
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation
from docx import Document

logger = logging.getLogger("study_notes.files")

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".txt",
                        ".xlsx", ".xlsm", ".html", ".htm"}


def extract_pdf_text(path: Path) -> str:
    try:
        reader = PdfReader(str(path))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"[Page {i + 1}]\n{text.strip()}")
        result = "\n\n".join(pages)
        logger.debug("PDF '%s': extracted %d chars from %d pages", path.name, len(result), len(reader.pages))
        return result
    except Exception as exc:
        logger.warning("Could not extract PDF '%s': %s", path.name, exc)
        return ""


def extract_pptx_text(path: Path) -> str:
    try:
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(r.text for r in para.runs).strip()
                        if line:
                            texts.append(line)
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        result = "\n\n".join(slides)
        logger.debug("PPTX '%s': extracted %d chars from %d slides", path.name, len(result), len(prs.slides))
        return result
    except Exception as exc:
        logger.warning("Could not extract PPTX '%s': %s", path.name, exc)
        return ""


def extract_docx_text(path: Path) -> str:
    try:
        doc = Document(str(path))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        result = "\n\n".join(paragraphs)
        logger.debug("DOCX '%s': extracted %d chars", path.name, len(result))
        return result
    except Exception as exc:
        logger.warning("Could not extract DOCX '%s': %s", path.name, exc)
        return ""


def extract_xlsx_text(path: Path, max_rows: int = 300) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        sheets = []
        for ws in wb.worksheets:
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows:
                    break
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                sheets.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows))
        wb.close()
        result = "\n\n".join(sheets)
        logger.debug("XLSX '%s': extracted %d chars", path.name, len(result))
        return result
    except Exception as exc:
        logger.warning("Could not extract XLSX '%s': %s", path.name, exc)
        return ""


def extract_html_text(path: Path) -> str:
    try:
        from lxml import html as lxml_html
        root = lxml_html.parse(str(path)).getroot()
        for bad in root.xpath("//script | //style"):
            bad.getparent().remove(bad)
        text = root.text_content()
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        result = "\n".join(lines)
        logger.debug("HTML '%s': extracted %d chars", path.name, len(result))
        return result
    except Exception as exc:
        logger.warning("Could not extract HTML '%s': %s", path.name, exc)
        return ""


def extract_text(path: Path) -> str:
    """Auto-detect file type and extract text. Returns empty string if unsupported."""
    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_pdf_text(path)
    if ext in (".pptx", ".ppt"):
        return extract_pptx_text(path)
    if ext in (".docx", ".doc"):
        return extract_docx_text(path)
    if ext in (".xlsx", ".xlsm"):
        return extract_xlsx_text(path)
    if ext in (".html", ".htm"):
        return extract_html_text(path)
    if ext == ".txt":
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return ""
    logger.debug("Unsupported file type: %s", ext)
    return ""


def summarise_materials(paths: list[Path], max_chars_per_file: int = 8000) -> list[dict]:
    """Extract text from each file and return a list of {name, text} dicts."""
    results = []
    for path in paths:
        if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        text = extract_text(path)
        if text:
            results.append({"name": path.name, "text": text[:max_chars_per_file]})
    return results
